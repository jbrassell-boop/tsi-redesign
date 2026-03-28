"""
North Floor Meeting — Daily PPTX Generator
Queries WinScopeNet SQL Server and populates the floor meeting template.

Usage:
    python tasks/floor-meeting/generate.py
    python tasks/floor-meeting/generate.py --date 2026-03-26   # specific date
"""

import json
import re
import sys
from copy import deepcopy
from datetime import datetime, timedelta
from pathlib import Path

import pyodbc
from pptx import Presentation
from lxml import etree

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------
SCRIPT_DIR = Path(__file__).parent
CONFIG = json.loads((SCRIPT_DIR / "config.json").read_text())
TEMPLATE = SCRIPT_DIR / "template.pptx"

# Parse --date argument or default to today
meeting_date = datetime.now()
if "--date" in sys.argv:
    idx = sys.argv.index("--date")
    meeting_date = datetime.strptime(sys.argv[idx + 1], "%Y-%m-%d")

yesterday = (meeting_date - timedelta(days=1)).strftime("%Y-%m-%d")
week_ago = (meeting_date - timedelta(days=7)).strftime("%Y-%m-%d")

# Location filter: north (default) or south
location = "north"
if "--location" in sys.argv:
    idx = sys.argv.index("--location")
    location = sys.argv[idx + 1].lower()

LOCATION_PREFIX = "NR" if location == "north" else "SR"
LOCATION_LABEL = "North" if location == "north" else "South"
NORTH_FILTER = f"AND r.sWorkOrderNumber LIKE '{LOCATION_PREFIX}%'"

# Skip internal/test clients flagged with bSkipTracking
SKIP_METRICS = """AND NOT EXISTS (
    SELECT 1 FROM tblDepartment d2
        JOIN tblClient c2 ON d2.lClientKey = c2.lClientKey
    WHERE d2.lDepartmentKey = r.lDepartmentKey
        AND c2.bSkipTracking = 1
)"""

# ---------------------------------------------------------------------------
# SQL Connection
# ---------------------------------------------------------------------------
_conn = None

def get_connection():
    global _conn
    if _conn is None:
        conn_str = (
            f"DRIVER={{ODBC Driver 18 for SQL Server}};"
            f"SERVER={CONFIG['sql_server']};"
            f"DATABASE={CONFIG['sql_database']};"
            f"Trusted_Connection=yes;"
            f"TrustServerCertificate=yes;"
        )
        _conn = pyodbc.connect(conn_str)
    return _conn

def query(sql, params=None):
    conn = get_connection()
    cursor = conn.cursor()
    cursor.execute(sql, params or [])
    columns = [col[0] for col in cursor.description]
    rows = [dict(zip(columns, row)) for row in cursor.fetchall()]
    return rows

# ---------------------------------------------------------------------------
# Scope type CASE expressions
# ---------------------------------------------------------------------------
SCOPE_CASE = """
    CASE
        WHEN st.sRigidOrFlexible = 'F' AND ISNULL(stc.bLargeDiameter, 0) = 1 THEN 'LD Flex'
        WHEN st.sRigidOrFlexible = 'F' THEN 'SD Flex'
        WHEN st.sRigidOrFlexible = 'R' THEN 'Rigid'
        WHEN st.sRigidOrFlexible = 'C' THEN 'Camera'
    END
"""

PRODUCTION_CASE = """
    CASE
        WHEN st.sRigidOrFlexible = 'F' AND ISNULL(stc.bLargeDiameter, 0) = 1 THEN 'LD'
        WHEN st.sRigidOrFlexible = 'F' THEN 'SD'
        WHEN st.sRigidOrFlexible = 'R' THEN 'Rigid'
        WHEN st.sRigidOrFlexible = 'C' THEN 'Cameras'
    END
"""

# Derive repair level from status (the status name encodes the level)
LEVEL_CASE = """
    CASE
        WHEN r.lRepairStatusID = 8  THEN 'Minor'
        WHEN r.lRepairStatusID = 11 THEN 'Mid'
        WHEN r.lRepairStatusID = 9  THEN 'Major'
        ELSE COALESCE(
            CASE
                WHEN rl.sRepairLevel = 'Minor'     THEN 'Minor'
                WHEN rl.sRepairLevel = 'Mid-Level'  THEN 'Mid'
                WHEN rl.sRepairLevel = 'Major'      THEN 'Major'
            END,
            '-'
        )
    END
"""

# Common scope join
SCOPE_JOIN = """
    JOIN tblScope s ON r.lScopeKey = s.lScopeKey
    JOIN tblScopeType st ON s.lScopeTypeKey = st.lScopeTypeKey
    LEFT JOIN tblScopeTypeCategories stc ON st.lScopeTypeCatKey = stc.lScopeTypeCategoryKey
"""

# In-house statuses (excludes outsourced, shipped, cart, more-info)
IN_HOUSE_STATUSES = "r.lRepairStatusID NOT IN (4, 10, 12, 13, 16, 17, 19, 20, 21, 22)"

# ---------------------------------------------------------------------------
# Queries
# ---------------------------------------------------------------------------

def query_slide3_flow():
    """Slide 3: Overall Flow — Previous Work Day"""
    received = query(f"""
        SELECT {SCOPE_CASE} AS ScopeCategory, COUNT(*) AS Cnt
        FROM tblRepair r {SCOPE_JOIN}
        WHERE CAST(r.dtDateIn AS DATE) = ?
            AND st.sRigidOrFlexible IN ('F','R','C')
            {NORTH_FILTER}
            {SKIP_METRICS}
        GROUP BY st.sRigidOrFlexible, ISNULL(stc.bLargeDiameter, 0)
    """, [yesterday])

    closed = query(f"""
        SELECT {SCOPE_CASE} AS ScopeCategory, COUNT(*) AS Cnt
        FROM tblRepair r {SCOPE_JOIN}
        WHERE CAST(r.dtDateOut AS DATE) = ?
            AND ISNULL(r.sRepairClosed, 'N') = 'Y'
            AND st.sRigidOrFlexible IN ('F','R','C')
            {NORTH_FILTER}
            {SKIP_METRICS}
        GROUP BY st.sRigidOrFlexible, ISNULL(stc.bLargeDiameter, 0)
    """, [yesterday])

    missed = query(f"""
        SELECT {SCOPE_CASE} AS ScopeCategory, COUNT(*) AS Cnt
        FROM tblRepair r {SCOPE_JOIN}
        WHERE CAST(r.dtExpDelDateTSI AS DATE) = ?
            AND ISNULL(r.sRepairClosed, 'N') <> 'Y'
            AND st.sRigidOrFlexible IN ('F','R','C')
            {NORTH_FILTER}
            {SKIP_METRICS}
        GROUP BY st.sRigidOrFlexible, ISNULL(stc.bLargeDiameter, 0)
    """, [yesterday])

    categories = ["LD Flex", "SD Flex", "Rigid", "Camera"]
    recv_map = {r["ScopeCategory"]: r["Cnt"] for r in received}
    closed_map = {r["ScopeCategory"]: r["Cnt"] for r in closed}
    missed_map = {r["ScopeCategory"]: r["Cnt"] for r in missed}

    rows = []
    for cat in categories:
        r_val = recv_map.get(cat, 0)
        c_val = closed_map.get(cat, 0)
        m_val = missed_map.get(cat, 0)
        net = r_val - c_val
        rows.append([cat, str(r_val), str(c_val), str(m_val), str(net)])

    totals = ["",
              str(sum(int(r[1]) for r in rows)),
              str(sum(int(r[2]) for r in rows)),
              str(sum(int(r[3]) for r in rows)),
              str(sum(int(r[4]) for r in rows))]
    rows.append(totals)
    return rows


def query_slide4_commitments():
    """Slide 4: Missed Commitments — Yesterday"""
    defects = query(f"""
        SELECT COUNT(DISTINCT dt.lRepairKey) AS Cnt
        FROM tblRepairDefectTracking dt
            JOIN tblRepair r ON dt.lRepairKey = r.lRepairKey
        WHERE CAST(r.dtDateIn AS DATE) = ?
            {NORTH_FILTER}
            {SKIP_METRICS}
    """, [yesterday])
    defect_count = defects[0]["Cnt"] if defects else 0

    loaner_data = query(f"""
        SELECT
            COUNT(*) AS Total,
            SUM(CASE WHEN lt.sDateOut IS NOT NULL AND lt.sDateOut <> '' THEN 1 ELSE 0 END) AS Fulfilled
        FROM tblLoanerTran lt
            JOIN tblRepair r ON lt.lRepairKey = r.lRepairKey
        WHERE ISNULL(r.sRepairClosed, 'N') <> 'Y'
            {NORTH_FILTER}
            {SKIP_METRICS}
    """)
    if loaner_data and loaner_data[0]["Total"] and loaner_data[0]["Total"] > 0:
        loaner_pct = round(100 * loaner_data[0]["Fulfilled"] / loaner_data[0]["Total"])
    else:
        loaner_pct = 0

    delivery_data = query(f"""
        SELECT
            COUNT(*) AS Total,
            SUM(CASE WHEN CAST(r.dtDateOut AS DATE) <= CAST(r.dtExpDelDateTSI AS DATE) THEN 1 ELSE 0 END) AS OnTime
        FROM tblRepair r
        WHERE CAST(r.dtDateOut AS DATE) = ?
            AND ISNULL(r.sRepairClosed, 'N') = 'Y'
            AND r.dtExpDelDateTSI IS NOT NULL
            {NORTH_FILTER}
            {SKIP_METRICS}
    """, [yesterday])
    if delivery_data and delivery_data[0]["Total"] and delivery_data[0]["Total"] > 0:
        delivery_pct = round(100 * delivery_data[0]["OnTime"] / delivery_data[0]["Total"])
    else:
        delivery_pct = 100

    kpi_rows = [
        ["Defect", str(defect_count), "<2"],
        ["Inspection Deadline", "Pass", "2pm"],
        ["Loaner Fulfillment", f"{loaner_pct}%", "75%"],
        ["Delivery Date", f"{delivery_pct}%", "92%"],
    ]

    missed_details = query(f"""
        SELECT
            {SCOPE_CASE} AS ScopeCategory,
            st.sScopeTypeDesc AS ScopeType,
            s.sSerialNumber AS SerialNumber,
            'Missed delivery' AS Reason,
            '' AS [Update]
        FROM tblRepair r {SCOPE_JOIN}
        WHERE CAST(r.dtExpDelDateTSI AS DATE) = ?
            AND ISNULL(r.sRepairClosed, 'N') <> 'Y'
            AND st.sRigidOrFlexible IN ('F','R','C')
            {NORTH_FILTER}
            {SKIP_METRICS}
    """, [yesterday])

    detail_rows = []
    for d in missed_details:
        detail_rows.append([
            d["ScopeCategory"] or "",
            d["ScopeType"] or "",
            d["SerialNumber"] or "",
            d["Reason"] or "",
            d["Update"] or "",
        ])
    if not detail_rows:
        detail_rows.append(["", "", "", "", ""])
    return kpi_rows, detail_rows


def query_slide5_production():
    """Slide 5: In-House Production Plan"""
    open_data = query(f"""
        SELECT
            {PRODUCTION_CASE} AS ScopeType,
            {LEVEL_CASE} AS RepairLevel,
            COUNT(*) AS WO_Open,
            SUM(CASE WHEN ISNULL(r.lTechnicianKey, 0) > 0 THEN 1 ELSE 0 END) AS OnBench
        FROM tblRepair r {SCOPE_JOIN}
            LEFT JOIN tblRepairRevenueAndExpenses rre ON r.lRepairKey = rre.lRepairKey
            LEFT JOIN tblRepairLevels rl ON rre.lRepairLevelKey = rl.lRepairLevelKey
        WHERE ISNULL(r.sRepairClosed, 'N') <> 'Y'
            AND st.sRigidOrFlexible IN ('F','R','C')
            AND {IN_HOUSE_STATUSES}
            {NORTH_FILTER}
            {SKIP_METRICS}
        GROUP BY st.sRigidOrFlexible, ISNULL(stc.bLargeDiameter, 0),
                 {LEVEL_CASE}
    """)

    weekly_out = query(f"""
        SELECT
            {PRODUCTION_CASE} AS ScopeType,
            {LEVEL_CASE} AS RepairLevel,
            COUNT(*) AS WeeklyOut
        FROM tblRepair r {SCOPE_JOIN}
            LEFT JOIN tblRepairRevenueAndExpenses rre ON r.lRepairKey = rre.lRepairKey
            LEFT JOIN tblRepairLevels rl ON rre.lRepairLevelKey = rl.lRepairLevelKey
        WHERE CAST(r.dtDateOut AS DATE) BETWEEN ? AND ?
            AND ISNULL(r.sRepairClosed, 'N') = 'Y'
            AND st.sRigidOrFlexible IN ('F','R','C')
            {NORTH_FILTER}
            {SKIP_METRICS}
        GROUP BY st.sRigidOrFlexible, ISNULL(stc.bLargeDiameter, 0),
                 {LEVEL_CASE}
    """, [week_ago, yesterday])

    open_map = {(r["ScopeType"], r["RepairLevel"]): r for r in open_data}
    weekly_map = {(r["ScopeType"], r["RepairLevel"]): r["WeeklyOut"] for r in weekly_out}
    goals = CONFIG["weekly_goals"]

    row_defs = [
        ("LD", "Minor", "LD_Minor"),
        ("LD", "Mid",   "LD_Mid"),
        ("LD", "Major", "LD_Major"),
        ("SD", "Minor", "SD_Minor"),
        ("SD", "Mid",   "SD_Mid"),
        ("SD", "Major", "SD_Major"),
        ("Rigid", "-",  "Rigid"),
        ("Cameras", "-","Cameras"),
    ]

    rows = []
    for scope_type, level, goal_key in row_defs:
        data = open_map.get((scope_type, level), {})
        wo_open = data.get("WO_Open", 0)
        on_bench = data.get("OnBench", 0)
        wk_out = weekly_map.get((scope_type, level), 0)
        goal = goals.get(goal_key, 0)
        gap = round(goal - wo_open)
        rows.append([
            scope_type, level,
            str(wo_open), str(on_bench), str(wk_out),
            str(round(goal)), str(gap),
        ])

    totals = [
        "", "",
        str(sum(int(r[2]) for r in rows)),
        str(sum(int(r[3]) for r in rows)),
        str(sum(int(r[4]) for r in rows)),
        str(sum(int(r[5]) for r in rows)),
        str(sum(int(r[6]) for r in rows)),
    ]
    rows.append(totals)
    return rows


def query_slide6_instruments():
    """Slide 6: Instruments Plan"""
    data = query(f"""
        SELECT
            CASE WHEN r.lRepairStatusID = 4 THEN 'Outsourced'
                 ELSE 'In House'
            END AS Channel,
            COUNT(*) AS WO_Open,
            SUM(CASE WHEN CAST(r.dtExpDelDateTSI AS DATE) <= DATEADD(day, 1, CAST(GETDATE() AS DATE))
                     THEN 1 ELSE 0 END) AS ExpectedSoon,
            SUM(CASE WHEN r.lRepairStatusID = 6 THEN 1 ELSE 0 END) AS WaitingQuotes
        FROM tblRepair r
            JOIN tblScope s ON r.lScopeKey = s.lScopeKey
            JOIN tblScopeType st ON s.lScopeTypeKey = st.lScopeTypeKey
        WHERE ISNULL(r.sRepairClosed, 'N') <> 'Y'
            AND (st.sRigidOrFlexible IN ('I','') OR st.sRigidOrFlexible IS NULL)
            {NORTH_FILTER}
            {SKIP_METRICS}
        GROUP BY CASE WHEN r.lRepairStatusID = 4 THEN 'Outsourced' ELSE 'In House' END
    """)

    channel_map = {r["Channel"]: r for r in data}
    rows = []
    for ch in ["In House", "Outsourced", "Van Service"]:
        d = channel_map.get(ch, {})
        rows.append([
            ch,
            str(d.get("WO_Open", 0)),
            str(d.get("ExpectedSoon", 0)),
            str(d.get("WaitingQuotes", 0)),
        ])
    return rows


def query_slide7_delays():
    """Slide 7: Delay Tracking — More Info Needed + Outsourced + More Time"""
    data = query(f"""
        SELECT
            CASE
                WHEN r.lRepairStatusID = 19 THEN 'More info'
                WHEN r.lRepairStatusID = 4  THEN 'Outsourced'
                WHEN r.lRepairStatusID = 5  THEN 'More time'
            END AS Category,
            c.sClientName1 AS ClientName,
            d.sShipState AS ShipState,
            r.sWorkOrderNumber AS WO,
            st.sScopeTypeDesc AS Model,
            s.sSerialNumber AS SN
        FROM tblRepair r
            JOIN tblScope s ON r.lScopeKey = s.lScopeKey
            JOIN tblScopeType st ON s.lScopeTypeKey = st.lScopeTypeKey
            JOIN tblDepartment d ON r.lDepartmentKey = d.lDepartmentKey
            JOIN tblClient c ON d.lClientKey = c.lClientKey
        WHERE ISNULL(r.sRepairClosed, 'N') <> 'Y'
            AND r.lRepairStatusID IN (19, 4, 5)
            AND ISNULL(c.bSkipTracking, 0) = 0
            {NORTH_FILTER}
        ORDER BY
            CASE r.lRepairStatusID
                WHEN 19 THEN 1
                WHEN 5  THEN 2
                WHEN 4  THEN 3
            END,
            r.dtDateIn ASC
    """)

    rows = []
    for d in data:
        client = d["ClientName"] or ""
        state = d["ShipState"] or ""
        # Only append state if the client name doesn't already end with it
        if state and not client.rstrip().endswith(f"- {state}"):
            client = f"{client} - {state}"
        rows.append([
            d["Category"] or "",
            client,
            d["WO"] or "",
            d["Model"] or "",
            d["SN"] or "",
            "",  # Status / Due Date — blank for manual fill
        ])
    if not rows:
        rows.append(["", "", "", "", "", ""])
    return rows


# ---------------------------------------------------------------------------
# PPTX Table Helpers
# ---------------------------------------------------------------------------

def populate_table(table, data_rows, header_rows=1):
    """Replace table data rows with data_rows, preserving header and cell formatting."""
    tbl = table._tbl
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

    tr_elements = tbl.findall(".//a:tr", nsmap)

    # Use first data row as format template
    if len(tr_elements) > header_rows:
        template_tr = deepcopy(tr_elements[header_rows])
    else:
        template_tr = deepcopy(tr_elements[-1])

    # Remove all existing data rows
    for tr in tr_elements[header_rows:]:
        tbl.remove(tr)

    # Add new rows
    for row_data in data_rows:
        new_tr = deepcopy(template_tr)
        tc_elements = new_tr.findall(".//a:tc", nsmap)
        for ci, cell_text in enumerate(row_data):
            if ci < len(tc_elements):
                tc = tc_elements[ci]
                runs = tc.findall(".//a:r", nsmap)
                if runs:
                    for run in runs:
                        t_elem = run.find("a:t", nsmap)
                        if t_elem is not None:
                            t_elem.text = ""
                    t_elem = runs[0].find("a:t", nsmap)
                    if t_elem is not None:
                        t_elem.text = str(cell_text)
                else:
                    ps = tc.findall(".//a:p", nsmap)
                    if ps:
                        txBody = tc.find("a:txBody", nsmap)
                        for p in ps[1:]:
                            txBody.remove(p)
                        p = ps[0]
                        for r in p.findall("a:r", nsmap):
                            p.remove(r)
                        r_elem = etree.SubElement(p, f"{{{nsmap['a']}}}r")
                        t_elem = etree.SubElement(r_elem, f"{{{nsmap['a']}}}t")
                        t_elem.text = str(cell_text)
        tbl.append(new_tr)


def find_table(slide, table_index=0):
    """Find the Nth table shape on a slide."""
    tables = [s for s in slide.shapes if s.has_table]
    if table_index < len(tables):
        return tables[table_index].table
    return None


def format_date(dt):
    """Format date as 'March 27th, 2026' with proper ordinal suffix."""
    day = dt.day
    if day in (1, 21, 31):
        suffix = "st"
    elif day in (2, 22):
        suffix = "nd"
    elif day in (3, 23):
        suffix = "rd"
    else:
        suffix = "th"
    return dt.strftime(f"%B {day}{suffix}, %Y")


def update_slide1_date(slide, dt):
    """Update the date in the Slide 1 title.
    Template has 3 runs: 'Staffing & Coverage Adjustments' | 'March 27th' | ', 2026'
    We replace run 1 with the new date and clear run 2.
    """
    date_str = format_date(dt)
    months = r"(January|February|March|April|May|June|July|August|September|October|November|December)"

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if not re.search(months, para.text):
                continue
            runs = list(para.runs)
            # Find the run with the month name and replace all date runs
            for i, run in enumerate(runs):
                if re.search(months, run.text):
                    run.text = date_str
                    # Clear all subsequent runs (e.g. ", 2026")
                    for j in range(i + 1, len(runs)):
                        runs[j].text = ""
                    return


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    print(f"Generating {LOCATION_LABEL} floor meeting for {meeting_date.strftime('%B %d, %Y')}")
    print(f"  Location: {LOCATION_LABEL} ({LOCATION_PREFIX}%)")
    print(f"  Yesterday: {yesterday}")
    print(f"  Week range: {week_ago} to {yesterday}")
    print()

    prs = Presentation(str(TEMPLATE))

    # --- Slide 1: Update date ---
    update_slide1_date(prs.slides[0], meeting_date)
    print("Slide 1: Date updated")

    # --- Slide 3: Overall Flow ---
    print("Querying Slide 3: Overall Flow...")
    flow_rows = query_slide3_flow()
    table3 = find_table(prs.slides[2])
    if table3:
        populate_table(table3, flow_rows)
        print(f"  {len(flow_rows)} rows")

    # --- Slide 4: Missed Commitments ---
    print("Querying Slide 4: Missed Commitments...")
    kpi_rows, detail_rows = query_slide4_commitments()
    table4a = find_table(prs.slides[3], 0)
    if table4a:
        populate_table(table4a, kpi_rows)
        print(f"  KPI: {len(kpi_rows)} rows")
    table4b = find_table(prs.slides[3], 1)
    if table4b:
        populate_table(table4b, detail_rows)
        print(f"  Details: {len(detail_rows)} rows")

    # --- Slide 5: Production Plan ---
    print("Querying Slide 5: In-House Production Plan...")
    prod_rows = query_slide5_production()
    table5 = find_table(prs.slides[4])
    if table5:
        populate_table(table5, prod_rows)
        print(f"  {len(prod_rows)} rows")

    # --- Slide 6: Instruments Plan ---
    print("Querying Slide 6: Instruments Plan...")
    inst_rows = query_slide6_instruments()
    table6 = find_table(prs.slides[5])
    if table6:
        populate_table(table6, inst_rows)
        print(f"  {len(inst_rows)} rows")

    # --- Slide 7: Delay Tracking ---
    print("Querying Slide 7: Delay Tracking...")
    delay_rows = query_slide7_delays()
    table7 = find_table(prs.slides[6])
    if table7:
        populate_table(table7, delay_rows)
        print(f"  {len(delay_rows)} rows")

    # --- Save ---
    output_dir = Path(CONFIG["output_dir"])
    filename = f"{LOCATION_LABEL} Floor Meeting.pptx"
    output_path = output_dir / filename
    prs.save(str(output_path))
    print(f"\nSaved to: {output_path}")
    print("Done!")

    # Close connection
    global _conn
    if _conn:
        _conn.close()


if __name__ == "__main__":
    main()
